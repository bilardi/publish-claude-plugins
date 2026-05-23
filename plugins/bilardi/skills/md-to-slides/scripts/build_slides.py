"""Rebuild the .pptx content slides from a Markdown source."""

import argparse
import re
import sys
from copy import deepcopy

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

import config


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

SOLAR_BG = RGBColor(*config.SOLAR_BG)
SOLAR_DEFAULT = RGBColor(*config.SOLAR_DEFAULT)
SOLAR_COMMENT = RGBColor(*config.SOLAR_COMMENT)
SOLAR_KEYWORD = RGBColor(*config.SOLAR_KEYWORD)
SOLAR_STRING = RGBColor(*config.SOLAR_STRING)
SOLAR_NUMBER = RGBColor(*config.SOLAR_NUMBER)
SOLAR_BUILTIN = RGBColor(*config.SOLAR_BUILTIN)


# ----- MD parser -----

def parse_md(md_text):
    slides = []
    current = None
    section = None
    in_code = False
    code_lang = None
    code_buf = []

    for raw in md_text.split("\n"):
        line = raw.rstrip()

        if in_code:
            if re.match(r"^```\s*$", line):
                target = section or "main"
                current["sections"][target] = {
                    "kind": "code",
                    "lang": code_lang,
                    "text": "\n".join(code_buf),
                }
                in_code = False
            else:
                code_buf.append(raw)
            continue

        m = re.match(r"^## type:(\d+)\s*(.*)$", line)
        if m:
            if current:
                slides.append(current)
            current = {"type": int(m.group(1)), "title": m.group(2).strip(), "sections": {}}
            section = None
            continue

        if current is None:
            continue

        m = re.match(r"^### (.+)$", line)
        if m:
            section = m.group(1).strip()
            current["sections"].setdefault(section, None)
            continue

        m = re.match(r"^```(\w*)$", line)
        if m:
            in_code = True
            code_lang = m.group(1) or "python"
            code_buf = []
            continue

        m = re.match(r"^(\s*)- (.+)$", line)
        if m:
            level = len(m.group(1)) // 2
            target = section or "main"
            data = current["sections"].get(target)
            if not (isinstance(data, dict) and data.get("kind") == "list"):
                data = {"kind": "list", "items": []}
                current["sections"][target] = data
            data["items"].append({"text": m.group(2), "level": level})
            continue

        m = re.match(r"^\|(.+)\|\s*$", line)
        if m and section:
            cells = [c.strip() for c in m.group(1).split("|")]
            # skip markdown table separator like |---|---|
            if all(re.match(r"^:?-+:?$", c) for c in cells if c):
                continue
            target = current["sections"].get(section)
            if not (isinstance(target, dict) and target.get("kind") == "table"):
                target = {"kind": "table", "rows": []}
                current["sections"][section] = target
            target["rows"].append(cells)
            continue

        if section and line.strip():
            data = current["sections"].get(section)
            if data is None:
                current["sections"][section] = {"kind": "text", "text": line}
            elif isinstance(data, dict) and data.get("kind") == "text":
                data["text"] += "\n" + line

    if current:
        slides.append(current)
    return slides


def parse_inline(text):
    chunks = []
    i = 0
    while i < len(text):
        m = re.match(r"\*\*\[([^\]]+)\]\(([^)]+)\)\*\*", text[i:])
        if m:
            chunks.append({"text": m.group(1), "bold": True, "link": m.group(2)})
            i += m.end()
            continue
        m = re.match(r"\*\*([^*]+)\*\*", text[i:])
        if m:
            chunks.append({"text": m.group(1), "bold": True, "link": None})
            i += m.end()
            continue
        m = re.match(r"\[([^\]]+)\]\(([^)]+)\)", text[i:])
        if m:
            chunks.append({"text": m.group(1), "bold": False, "link": m.group(2)})
            i += m.end()
            continue
        end = len(text)
        for marker in ("**", "["):
            pos = text.find(marker, i)
            if 0 <= pos < end:
                end = pos
        if end == i:
            end = i + 1
        chunks.append({"text": text[i:end], "bold": False, "link": None})
        i = end
    merged = []
    for c in chunks:
        if (
            merged
            and not c["bold"] and not c["link"]
            and not merged[-1]["bold"] and not merged[-1]["link"]
        ):
            merged[-1]["text"] += c["text"]
        else:
            merged.append(c)
    return merged


# ----- tokenizer -----

_PY_TOK = re.compile(
    r'(?P<comment>#[^\n]*)'
    r'|(?P<string>f?"[^"]*"|f?\'[^\']*\')'
    r'|(?P<number>\b\d+\b)'
    r'|(?P<word>\b[A-Za-z_]\w*\b)'
    r'|(?P<other>.)',
    re.DOTALL,
)


def tokenize_line(line, lang):
    if lang in ("shell", "bash", "console"):
        if line.startswith("$ "):
            yield SOLAR_KEYWORD, "$ "
            if line[2:]:
                yield SOLAR_DEFAULT, line[2:]
            return
        if line.startswith("#"):
            yield SOLAR_COMMENT, line
            return
        yield SOLAR_DEFAULT, line
        return
    for m in _PY_TOK.finditer(line):
        kind, text = m.lastgroup, m.group()
        if kind == "word":
            if text in config.PY_KEYWORDS:
                yield SOLAR_KEYWORD, text
            elif text in config.PY_BUILTINS:
                yield SOLAR_BUILTIN, text
            else:
                yield SOLAR_DEFAULT, text
        elif kind == "comment":
            yield SOLAR_COMMENT, text
        elif kind == "string":
            yield SOLAR_STRING, text
        elif kind == "number":
            yield SOLAR_NUMBER, text
        else:
            yield SOLAR_DEFAULT, text


# ----- slide helpers -----

def is_title(shape):
    try:
        return shape.placeholder_format.idx == 0 and shape.has_text_frame
    except (AttributeError, ValueError):
        return False


def is_code_box(shape):
    try:
        return shape.fill.type is not None and str(shape.fill.fore_color.rgb) == "002B36"
    except Exception:
        return False


def get_content_shapes(slide):
    out = [s for s in slide.shapes if s.has_text_frame and not is_title(s)]
    out.sort(key=lambda s: s.left or 0)
    return out


def find_markers(prs):
    out = []
    for i, slide in enumerate(prs.slides):
        if slide.slide_layout.name != "BLANK":
            continue
        has_text = any(
            sh.has_text_frame and sh.text_frame.text.strip()
            for sh in slide.shapes
        )
        if not has_text:
            out.append(i)
    return out


# ----- clone & replace -----

def clone_slide(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh.element.getparent().remove(sh.element)
    for sh in src.shapes:
        new.shapes._spTree.append(deepcopy(sh.element))
    src_bg = src.element.find(f"{{{P_NS}}}cSld/{{{P_NS}}}bg")
    if src_bg is not None:
        bg_copy = deepcopy(src_bg)
        for blip in bg_copy.iter(f"{{{A_NS}}}blip"):
            old_rid = blip.get(f"{{{R_NS}}}embed")
            if old_rid:
                image_part = src.part.related_part(old_rid)
                new_rid = new.part.relate_to(image_part, IMG_REL)
                blip.set(f"{{{R_NS}}}embed", new_rid)
        new.element.find(f"{{{P_NS}}}cSld").insert(0, bg_copy)
    return new


def _set_text_in_first_run(tf, text):
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    else:
        tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = config.TITLE_FONT


def set_title(slide, text):
    for sh in slide.shapes:
        if is_title(sh):
            _set_text_in_first_run(sh.text_frame, text)
            return


def set_subtitle(slide, text):
    boxes = [s for s in slide.shapes if s.has_text_frame and not is_title(s)]
    if boxes:
        _set_text_in_first_run(boxes[0].text_frame, text)


def replace_list(shape, items):
    tf = shape.text_frame
    src_paragraphs = list(tf.paragraphs)
    if not src_paragraphs:
        return
    # collect template paragraphs grouped by level, preserving order
    templates_by_level = {}
    for p in src_paragraphs:
        pPr = p._p.find(f"{{{A_NS}}}pPr")
        lvl = int(pPr.get("lvl", "0")) if pPr is not None else 0
        templates_by_level.setdefault(lvl, []).append(deepcopy(p._p))
    if 0 not in templates_by_level:
        templates_by_level[0] = [deepcopy(src_paragraphs[0]._p)]

    txBody = tf._txBody
    for p in src_paragraphs:
        p._p.getparent().remove(p._p)

    bullet_chars = ["●", "○", "▪"]
    count_by_level = {}
    for item in items:
        lvl = item["level"]
        if lvl in templates_by_level:
            templates = templates_by_level[lvl]
            do_fallback_override = False
        else:
            templates = templates_by_level[0]
            do_fallback_override = True
        idx = count_by_level.get(lvl, 0)
        new_p = deepcopy(templates[min(idx, len(templates) - 1)])
        count_by_level[lvl] = idx + 1

        if do_fallback_override:
            pPr = new_p.find(f"{{{A_NS}}}pPr")
            if pPr is not None:
                pPr.set("lvl", str(lvl))
                pPr.set("marL", str(457200 * (lvl + 1)))
                buChar = pPr.find(f"{{{A_NS}}}buChar")
                if buChar is not None:
                    buChar.set("char", bullet_chars[min(lvl, len(bullet_chars) - 1)])

        # extract endParaRPr to re-append after runs (correct schema order)
        endParaRPr = new_p.find(f"{{{A_NS}}}endParaRPr")
        if endParaRPr is not None:
            new_p.remove(endParaRPr)

        runs = new_p.findall(f"{{{A_NS}}}r")
        run_template = deepcopy(runs[0]) if runs else None
        for r_el in runs:
            new_p.remove(r_el)

        for chunk in parse_inline(item["text"]):
            if run_template is not None:
                r = deepcopy(run_template)
            else:
                r = etree.SubElement(etree.Element("dummy"), f"{{{A_NS}}}r")
                etree.SubElement(r, f"{{{A_NS}}}t").text = ""
            t = r.find(f"{{{A_NS}}}t")
            if t is not None:
                t.text = chunk["text"]
            rPr = r.find(f"{{{A_NS}}}rPr")
            if rPr is not None:
                had_link = rPr.find(f"{{{A_NS}}}hlinkClick") is not None
                for hlink in rPr.findall(f"{{{A_NS}}}hlinkClick"):
                    rPr.remove(hlink)
                if chunk.get("bold"):
                    rPr.set("b", "1")
                elif "b" in rPr.attrib:
                    del rPr.attrib["b"]
                if chunk.get("link"):
                    rPr.set("u", "sng")
                elif had_link:
                    rPr.set("u", "none")
                # else: leave u alone (preserves template default)
            new_p.append(r)
        if endParaRPr is not None:
            new_p.append(endParaRPr)
        txBody.append(new_p)

    new_paragraphs = list(tf.paragraphs)
    for i, item in enumerate(items):
        if i >= len(new_paragraphs):
            break
        chunks = parse_inline(item["text"])
        for j, run in enumerate(new_paragraphs[i].runs):
            if j < len(chunks) and chunks[j].get("link"):
                try:
                    run.hyperlink.address = chunks[j]["link"]
                except Exception:
                    pass


def replace_table(slide, rows_data):
    """Replace the table shape's content with rows_data (list of list of str).
    First row is the header, rest are data. Expands rows/columns when MD has more
    than the template; re-centers horizontally based on total column width."""
    if not rows_data:
        return
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl_xml = shape.table._tbl
        target_cols = max(len(r) for r in rows_data)
        target_rows = len(rows_data)

        while len(shape.table.columns) < target_cols:
            _add_table_column(tbl_xml)
        while len(shape.table.rows) < target_rows:
            _add_table_row(tbl_xml)

        for r_idx, row_data in enumerate(rows_data):
            for c_idx, cell_text in enumerate(row_data):
                cell = shape.table.cell(r_idx, c_idx)
                tf = cell.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = cell_text
                    for r in tf.paragraphs[0].runs[1:]:
                        r.text = ""
                else:
                    tf.text = cell_text

        total_w = sum(c.width for c in shape.table.columns)
        slide_width = slide.part.package.presentation_part.presentation.slide_width
        shape.left = (slide_width - total_w) // 2
        return


def _add_table_column(tbl_xml):
    """Append a column to the table by cloning the last gridCol and the last tc
    in each row (per-row, so header keeps header style and data keeps data style)."""
    grid = tbl_xml.find(f"{{{A_NS}}}tblGrid")
    if grid is None:
        return
    grid_cols = grid.findall(f"{{{A_NS}}}gridCol")
    if not grid_cols:
        return
    grid.append(deepcopy(grid_cols[-1]))
    for tr in tbl_xml.findall(f"{{{A_NS}}}tr"):
        tcs = tr.findall(f"{{{A_NS}}}tc")
        if not tcs:
            continue
        new_tc = deepcopy(tcs[-1])
        _clear_tc_text(new_tc)
        tr.append(new_tc)


def _add_table_row(tbl_xml):
    """Append a row by cloning the last existing row (data row, not header)."""
    trs = tbl_xml.findall(f"{{{A_NS}}}tr")
    if not trs:
        return
    new_tr = deepcopy(trs[-1])
    for tc in new_tr.findall(f"{{{A_NS}}}tc"):
        _clear_tc_text(tc)
    tbl_xml.append(new_tr)


def _clear_tc_text(tc_el):
    """Empty all runs' text in a <a:tc> while preserving formatting."""
    txBody = tc_el.find(f"{{{A_NS}}}txBody")
    if txBody is None:
        return
    for p in txBody.findall(f"{{{A_NS}}}p"):
        for r in p.findall(f"{{{A_NS}}}r"):
            t = r.find(f"{{{A_NS}}}t")
            if t is not None:
                t.text = ""


def replace_picture(slide, image_path, ref_width=None):
    """Swap image, centered horizontally on the slide.

    - If `ref_width` is given: use it as the width; compute height proportional
      to the image; keep template's top.
    - Else: use template's top + height; compute width proportional to the image."""
    from PIL import Image as PILImage
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        top = shape.top
        template_height = shape.height
        shape.element.getparent().remove(shape.element)
        with PILImage.open(image_path) as img:
            iw, ih = img.size
        if ref_width is not None:
            width = ref_width
            height = int(width * ih / iw)
        else:
            height = template_height
            width = int(template_height * iw / ih)
        slide_width = slide.part.package.presentation_part.presentation.slide_width
        left = (slide_width - width) // 2
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        return


def compute_type7_ref_width(prs, slides_def):
    """Compute the reference width for type:7 images using the last one referenced.
    Width = template_height * (image_width / image_height). Returns None if no type:7
    image is referenced."""
    from PIL import Image as PILImage
    image_paths = []
    for s in slides_def:
        if s["type"] != 7:
            continue
        data = s.get("sections", {}).get("image")
        if isinstance(data, dict):
            path = data.get("text", "").strip()
            if path:
                image_paths.append(path)
    if not image_paths:
        return None
    template_pos = config.TYPE_TO_TEMPLATE.get(7)
    if template_pos is None:
        return None
    template_slide = prs.slides[template_pos - 1]
    template_height = None
    for shape in template_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            template_height = shape.height
            break
    if template_height is None:
        return None
    with PILImage.open(image_paths[-1]) as img:
        iw, ih = img.size
    return int(template_height * iw / ih)


def replace_code(shape, code_text, lang):
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = SOLAR_BG
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()
    for li, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        if not line:
            r = p.add_run()
            r.text = ""
            _format_run(r, SOLAR_DEFAULT)
            continue
        for color, text in tokenize_line(line, lang):
            r = p.add_run()
            r.text = text
            _format_run(r, color)


def _format_run(run, color):
    run.font.name = config.TITLE_FONT
    run.font.size = Pt(config.TEXT_PT)
    run.font.color.rgb = color


# ----- per-type content apply -----

def apply_content(slide, slide_def, type7_ref_width=None):
    set_title(slide, slide_def["title"])
    typ = slide_def["type"]
    sections = slide_def.get("sections", {})

    if typ in (3, 5):
        if isinstance(sections.get("subtitle"), dict):
            set_subtitle(slide, sections["subtitle"].get("text", ""))
        return

    if typ == 2:
        boxes = get_content_shapes(slide)
        for sec in ("items", "main"):
            data = sections.get(sec)
            if isinstance(data, dict) and data.get("kind") == "list" and boxes:
                replace_list(boxes[0], data["items"])
                return
        return

    if typ in (4, 6):
        boxes = get_content_shapes(slide)
        if len(boxes) < 2:
            return
        for side, box in (("left", boxes[0]), ("right", boxes[1])):
            data = sections.get(side)
            if not isinstance(data, dict):
                continue
            if data["kind"] == "list":
                replace_list(box, data["items"])
            elif data["kind"] == "code":
                replace_code(box, data["text"], data.get("lang") or "python")
        return

    if typ == 7:
        data = sections.get("image")
        if isinstance(data, dict):
            img_path = data.get("text", "").strip()
            if img_path:
                replace_picture(slide, img_path, ref_width=type7_ref_width)
        return

    if typ == 8:
        data = sections.get("table")
        if isinstance(data, dict) and data.get("kind") == "table":
            replace_table(slide, data["rows"])
        return


# ----- build -----

def build(slides_def, pptx_path):
    """Two-pass build: clear+save (drops orphans), then reopen and add."""
    prs = Presentation(pptx_path)

    # type:2 (Agenda) is a fixed-position in-place update, not cloned
    agenda_defs = [s for s in slides_def if s["type"] == 2]
    if agenda_defs:
        apply_content(prs.slides[config.AGENDA_POSITION - 1], agenda_defs[0])
    other_defs = [s for s in slides_def if s["type"] != 2]

    markers = find_markers(prs)
    if len(markers) < 2:
        sys.exit(f"Need 2 BLANK marker slides; found {len(markers)}")
    start_idx, end_idx = markers[0], markers[1]

    sld_id_lst = prs.slides._sldIdLst
    rid_attr = f"{{{R_NS}}}id"
    for el in list(sld_id_lst)[start_idx + 1 : end_idx]:
        rId = el.get(rid_attr)
        sld_id_lst.remove(el)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
    prs.save(pptx_path)

    # reopen so orphan slide parts are flushed
    prs = Presentation(pptx_path)
    sld_id_lst = prs.slides._sldIdLst
    type7_ref_width = compute_type7_ref_width(prs, other_defs)

    for slide_def in other_defs:
        typ = slide_def["type"]
        template_pos = config.TYPE_TO_TEMPLATE.get(typ)
        if template_pos is None:
            print(f"  skip unknown type:{typ}", file=sys.stderr)
            continue
        template_slide = prs.slides[template_pos - 1]
        new_slide = clone_slide(prs, template_slide)
        apply_content(new_slide, slide_def, type7_ref_width=type7_ref_width)

        children = list(sld_id_lst)
        new_id = children[-1]
        sld_id_lst.remove(new_id)
        markers = find_markers(prs)
        if len(markers) < 2:
            sys.exit("Lost markers during build")
        sld_id_lst.insert(markers[1], new_id)

    prs.save(pptx_path)
    return prs


def main():
    parser = argparse.ArgumentParser(description="Build .pptx content slides from .md")
    parser.add_argument("--md", default=config.DEFAULT_MD, help="input MD path")
    parser.add_argument("--pptx", default=config.DEFAULT_PPTX, help="output .pptx path")
    args = parser.parse_args()

    with open(args.md) as f:
        slides_def = parse_md(f.read())
    prs = build(slides_def, args.pptx)
    print(f"Built {args.pptx}: {len(prs.slides)} slides total")


if __name__ == "__main__":
    main()
