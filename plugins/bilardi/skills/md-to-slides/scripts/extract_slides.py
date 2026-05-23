"""Read the .pptx between markers and write a Markdown source file."""

import argparse
import sys

from pptx import Presentation

import config


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def is_title(shape):
    try:
        return shape.placeholder_format.idx == 0 and shape.has_text_frame
    except (AttributeError, ValueError):
        return False


def is_code_box(shape):
    try:
        return shape.fill.type is not None and str(shape.fill.fore_color.rgb) == "002B36"
    except Exception:
        return False


def slide_title(slide):
    for s in slide.shapes:
        if is_title(s):
            return s.text_frame.text
    return ""


def get_content_shapes(slide):
    out = [s for s in slide.shapes if s.has_text_frame and not is_title(s)]
    out.sort(key=lambda s: s.left or 0)
    return out


def find_markers(prs):
    out = []
    for i, slide in enumerate(prs.slides):
        if slide.slide_layout.name != "BLANK":
            continue
        has_text = any(
            sh.has_text_frame and sh.text_frame.text.strip()
            for sh in slide.shapes
        )
        if not has_text:
            out.append(i)
    return out


def detect_type(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    layout = slide.slide_layout.name
    has_bg = slide.element.find(f"{{{P_NS}}}cSld/{{{P_NS}}}bg") is not None
    if layout == "TITLE":
        return 3 if has_bg else 5
    if layout == "TITLE_AND_BODY":
        boxes = get_content_shapes(slide)
        if len(boxes) >= 2:
            l, r = is_code_box(boxes[0]), is_code_box(boxes[1])
            if l and r:
                return 6
            if not l and r:
                return 4
    if layout == "BIG_NUMBER":
        return 2
    if layout == "TITLE_ONLY_1":
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return 8
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return 7
    return None


def detect_lang(box):
    text = box.text_frame.text
    if text.startswith("$") or "\n$ " in text:
        return "shell"
    return "python"


def render_list_item(paragraph):
    pPr = paragraph._p.find(f"{{{A_NS}}}pPr")
    lvl = int(pPr.get("lvl", "0")) if pPr is not None else 0
    indent = "  " * lvl
    chunks = []
    for r in paragraph.runs:
        text = r.text or ""
        if not text:
            continue
        bold = r.font.bold
        link = None
        try:
            if r.hyperlink and r.hyperlink.address:
                link = r.hyperlink.address
        except Exception:
            pass
        if bold and link:
            chunks.append(f"**[{text}]({link})**")
        elif bold:
            chunks.append(f"**{text}**")
        elif link:
            chunks.append(f"[{text}]({link})")
        else:
            chunks.append(text)
    full = "".join(chunks).strip()
    if not full:
        return None
    return f"{indent}- {full}"


def extract(prs):
    markers = find_markers(prs)
    if len(markers) < 2:
        sys.exit("Need 2 BLANK markers")
    start_idx, end_idx = markers[0], markers[1]

    lines = ["# Workshop slides", ""]

    # Agenda (fixed position, before markers)
    agenda_slide = prs.slides[config.AGENDA_POSITION - 1]
    lines.append(f"## type:2 {slide_title(agenda_slide)}")
    boxes = get_content_shapes(agenda_slide)
    if boxes:
        lines.append("### items")
        for p in boxes[0].text_frame.paragraphs:
            rendered = render_list_item(p)
            if rendered is not None:
                lines.append(rendered)
    lines.append("")

    for i in range(start_idx + 1, end_idx):
        slide = prs.slides[i]
        typ = detect_type(slide)
        title = slide_title(slide)
        if typ is None:
            lines.append(f"## UNKNOWN slide pos {i + 1}: {title}")
            lines.append("")
            continue
        lines.append(f"## type:{typ} {title}")

        if typ in (3, 5):
            boxes = [s for s in slide.shapes if s.has_text_frame and not is_title(s)]
            if boxes:
                subtitle = boxes[0].text_frame.text
                if subtitle.strip():
                    lines.append("### subtitle")
                    lines.append(subtitle)
        elif typ in (4, 6):
            boxes = get_content_shapes(slide)
            for side, box in zip(("left", "right"), boxes[:2]):
                if is_code_box(box):
                    lang = detect_lang(box)
                    text = "\n".join(
                        "".join(r.text for r in p.runs)
                        for p in box.text_frame.paragraphs
                    )
                    lines.append(f"### {side}")
                    lines.append(f"```{lang}")
                    lines.append(text)
                    lines.append("```")
                else:
                    lines.append(f"### {side}")
                    for p in box.text_frame.paragraphs:
                        rendered = render_list_item(p)
                        if rendered is not None:
                            lines.append(rendered)
        elif typ == 2:
            boxes = get_content_shapes(slide)
            if boxes:
                lines.append("### items")
                for p in boxes[0].text_frame.paragraphs:
                    rendered = render_list_item(p)
                    if rendered is not None:
                        lines.append(rendered)
        elif typ == 7:
            # original image path is not stored in the .pptx; emit placeholder
            lines.append("### image")
            lines.append("images/?.png  # TODO: set actual path")
        elif typ == 8:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                    continue
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text_frame.text for cell in row.cells])
                if rows:
                    lines.append("### table")
                    lines.append("| " + " | ".join(rows[0]) + " |")
                    lines.append("|" + "|".join(["---"] * len(rows[0])) + "|")
                    for row in rows[1:]:
                        lines.append("| " + " | ".join(row) + " |")
                break
        lines.append("")
    return "\n".join(lines) + "\n"


def main():
    parser = argparse.ArgumentParser(description="Extract slide content from .pptx to .md")
    parser.add_argument("--md", default=config.DEFAULT_MD, help="output MD path")
    parser.add_argument("--pptx", default=config.DEFAULT_PPTX, help="input .pptx path")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    md = extract(prs)
    with open(args.md, "w") as f:
        f.write(md)
    print(f"Extracted {args.md}")


if __name__ == "__main__":
    main()
